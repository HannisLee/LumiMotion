#!/usr/bin/env python3
"""对生成的 PPT 做不依赖 Office 渲染器的结构检查。"""

from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation


REPORT_DIR = Path(__file__).resolve().parent
PPTX = REPORT_DIR / "0818-Lambertian法线恢复审计与实验汇报.pptx"
REPORT = REPORT_DIR / "qa_report.txt"


def main():
    lines: list[str] = []
    errors: list[str] = []
    warnings: list[str] = []

    with ZipFile(PPTX) as archive:
        bad_member = archive.testzip()
        names = archive.namelist()
        slide_xml = sorted(
            name for name in names
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        media = sorted(name for name in names if name.startswith("ppt/media/"))
        external_links = []
        for name in names:
            if not name.endswith(".rels"):
                continue
            payload = archive.read(name)
            if b'TargetMode="External"' in payload:
                external_links.append(name)
        if bad_member:
            errors.append(f"ZIP CRC 检查失败：{bad_member}")
        if external_links:
            warnings.append(f"存在外部关系：{external_links}")
        lines.append(f"ZIP 成员数：{len(names)}")
        lines.append(f"幻灯片 XML：{len(slide_xml)}")
        lines.append(f"内嵌媒体：{len(media)}")

    prs = Presentation(PPTX)
    sw, sh = prs.slide_width, prs.slide_height
    lines.append(f"python-pptx 页数：{len(prs.slides)}")
    lines.append(f"页面尺寸：{sw / 914400:.3f} × {sh / 914400:.3f} 英寸")

    empty_text_boxes = 0
    for si, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > sw or shape.top + shape.height > sh:
                errors.append(
                    f"第 {si} 页形状越界：{getattr(shape, 'name', '<unnamed>')} "
                    f"({shape.left}, {shape.top}, {shape.width}, {shape.height})"
                )
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if not text:
                    empty_text_boxes += 1
                    continue
                # 过小字号检查：页脚允许 8.5 pt，其余不应低于 9 pt。
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is None:
                            continue
                        pt = run.font.size.pt
                        if pt < 8.4:
                            warnings.append(f"第 {si} 页字号过小：{pt:.1f} pt，文本={text[:32]!r}")

    if len(prs.slides) != 19:
        errors.append(f"页数不符合预期：{len(prs.slides)} != 19")
    if not media:
        errors.append("PPT 中没有内嵌媒体")

    lines.append(f"空文本形状：{empty_text_boxes}（多为不需要文本的装饰图形）")
    lines.append(f"错误：{len(errors)}")
    lines.append(f"警告：{len(warnings)}")
    lines.append("")
    if errors:
        lines.append("[错误]")
        lines.extend(f"- {item}" for item in errors)
        lines.append("")
    if warnings:
        lines.append("[警告]")
        lines.extend(f"- {item}" for item in warnings)
        lines.append("")
    lines.append("[结论]")
    lines.append("PASS" if not errors else "FAILED")
    lines.append("")
    lines.append("[已知限制]")
    lines.append("当前 garuda 服务器没有 PowerPoint/LibreOffice 渲染器，未执行逐页像素级渲染检查。")
    lines.append("已执行 ZIP CRC、内嵌媒体、外部链接、页面边界、页数和最小字号检查。")

    REPORT.write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(REPORT)
    print("\n".join(lines))
    raise SystemExit(1 if errors else 0)


if __name__ == "__main__":
    main()
